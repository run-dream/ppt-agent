import os
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Cm
from PIL import Image
from src.models.state import PPTState, PPTOutline, SlideContent
from src.utils.layout_manager import LayoutManager
from src.utils.logger import logger

class PPTGenerator:
    """
    PPT 生成器：根据 PPTState 中的数据渲染并导出 .pptx 文件
    """
    
    def __init__(self, output_dir: str = "data/outputs"):
        self.output_dir = output_dir
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

    def generate(self, state: PPTState) -> str:
        """
        根据状态生成 PPT 文件并返回文件路径
        """
        outline = state.get("outline")
        slides_data = state.get("slides", [])

        logger.info(f"PPT Generator: Outline title: '{outline.title if outline else 'None'}'")
        logger.info(f"PPT Generator: Number of slides: {len(slides_data)}")

        if not outline:
            raise ValueError("No outline found in state for PPT generation.")

        prs = Presentation()

        # 检查是否有标题页数据
        title_slide_data = None
        content_slides_data = slides_data

        if slides_data and len(slides_data) > 0 and getattr(slides_data[0], 'layout_type', None) == 'title_slide':
            title_slide_data = slides_data[0]
            content_slides_data = slides_data[1:]  # 剩余的是内容页
            logger.info("PPT Generator: Found title slide data from generator")

        # 1. 创建标题页
        if title_slide_data:
            # 使用generator生成的标题页数据
            cover_image = getattr(title_slide_data, 'image_path', None)
            self._add_title_slide(prs, title_slide_data.title, cover_image)
            # 如果标题页有bullet_points，也添加到标题页
            if hasattr(title_slide_data, 'bullet_points') and title_slide_data.bullet_points:
                # 这里可以扩展_add_title_slide方法来添加副标题
                logger.info(f"PPT Generator: Title slide subtitle: {title_slide_data.bullet_points}")
        else:
            # 回退到原来的逻辑
            cover_image = None
            if slides_data and hasattr(slides_data[0], 'image_path') and slides_data[0].image_path:
                cover_image = slides_data[0].image_path
            self._add_title_slide(prs, outline.title, cover_image)

        logger.info("PPT Generator: Title slide added")

        # 2. 逐页创建内容幻灯片
        for slide_data in content_slides_data:
            self._add_content_slide(prs, slide_data)
            
        # 3. 保存文件
        safe_title = "".join([c for c in outline.title if c.isalnum() or c in (' ', '_')]).rstrip()
        filename = f"{safe_title or 'presentation'}.pptx"
        file_path = os.path.join(self.output_dir, filename)
        
        prs.save(file_path)
        logger.info(f"PPT successfully generated at: {file_path}")
        logger.info(f"PPT contains {len(prs.slides)} slides")
        for i, slide in enumerate(prs.slides):
            logger.info(f"Slide {i+1}: {slide.slide_layout.name}")
        return file_path

    def _add_title_slide(self, prs, title_text: str, cover_image: str = None):
        logger.info(f"Adding title slide with title: '{title_text}'")

        layout = prs.slide_layouts[LayoutManager.get_layout_index("title_slide")]
        slide = prs.slides.add_slide(layout)
        logger.debug(f"Title slide created with layout index: {LayoutManager.get_layout_index('title_slide')}")

        # 设置标题
        title_ph = LayoutManager.get_placeholder(slide, 'title')
        if title_ph:
            title_ph.text = title_text
            logger.debug(f"Title text set: '{title_text}'")
        else:
            logger.warning("No title placeholder found in title slide layout")

        # 添加封面图片（标题页通常没有图片占位符，使用fallback方法）
        if cover_image:
            logger.info(f"Adding cover image to title slide: {cover_image}")
            # 标题页直接使用fallback方法添加图片，因为layout 0通常没有图片占位符
            self._add_image_force(slide, cover_image)

    def _add_content_slide(self, prs, data: SlideContent):
        layout_idx = LayoutManager.get_layout_index(data.layout_type)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        logger.debug(f"Content slide layout: {layout.name} (index: {layout_idx})")
        logger.debug(f"Available placeholders: {[ph.placeholder_format.type for ph in slide.placeholders]}")

        # 填充标题
        title_ph = LayoutManager.get_placeholder(slide, 'title')
        if title_ph:
            title_ph.text = data.title
        else:
            logger.warning(f"No title placeholder found in layout {layout_idx}")

        # 填充正文
        body_ph = LayoutManager.get_placeholder(slide, 'body')
        if body_ph:
            tf = body_ph.text_frame
            tf.clear() # 清除默认占位符文本
            for point in data.bullet_points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        else:
            logger.warning(f"No body placeholder found in layout {layout_idx}")

        # 填充图片 (如果有路径或 URL)
        if data.image_path:
            logger.debug(f"Adding image to content slide: {data.image_path}")
            self._add_image(slide, data.image_path)
        else:
            logger.debug("No image_path for content slide")

    def _add_image(self, slide, image_source: str):
        """添加图片到幻灯片，支持本地路径或 URL"""
        pic_ph = LayoutManager.get_placeholder(slide, 'picture')

        try:
            image_data = None
            if image_source.startswith(('http://', 'https://')):
                # 处理远程图片
                response = requests.get(image_source, timeout=10, stream=True)
                response.raise_for_status()
                image_data = BytesIO(response.content)

                # 检查是否是占位图URL，如果是则跳过
                if 'placehold.co' in image_source or 'placeholder.com' in image_source:
                    logger.info("Skipping placeholder image insertion")
                    return
            elif os.path.exists(image_source):
                # 处理本地图片
                with open(image_source, 'rb') as f:
                    image_data = BytesIO(f.read())
            else:
                logger.warning(f"Image source not found: {image_source}")
                return

            if not image_data:
                return

            # 尝试使用占位符插入图片
            if pic_ph:
                try:
                    pic_ph.insert_picture(image_data)
                    logger.debug("Image inserted using placeholder")
                    return
                except Exception as e:
                    logger.warning(f"Failed to insert image using placeholder: {e}")

            # Fallback: 直接在slide上添加图片
            self._add_image_force(slide, image_data)

        except Exception as e:
            logger.error(f"Failed to add image to slide: {str(e)}")
            logger.debug(f"Image source: {image_source}")

    def _add_image_force(self, slide, image_data_or_source: str):
        """强制添加图片到slide，使用fallback方法，不尝试占位符"""
        try:
            # 处理输入参数
            if isinstance(image_data_or_source, str):
                # 如果是字符串，当作URL或文件路径处理
                if image_data_or_source.startswith(('http://', 'https://')):
                    response = requests.get(image_data_or_source, timeout=10, stream=True)
                    response.raise_for_status()
                    image_data = BytesIO(response.content)
                elif os.path.exists(image_data_or_source):
                    with open(image_data_or_source, 'rb') as f:
                        image_data = BytesIO(f.read())
                else:
                    logger.warning(f"Image source not found: {image_data_or_source}")
                    return
            else:
                # 假设已经是BytesIO对象
                image_data = image_data_or_source

            # 获取图片尺寸
            image_data.seek(0)
            img = Image.open(image_data)
            img_width, img_height = img.size

            # 计算合适的显示尺寸 (保持宽高比，最大宽度8cm，最大高度6cm)
            max_width = Cm(8)
            max_height = Cm(6)
            ratio = min(max_width / Cm(2.54 * img_width / 96), max_height / Cm(2.54 * img_height / 96))
            display_width = int(Cm(2.54 * img_width / 96) * ratio)
            display_height = int(Cm(2.54 * img_height / 96) * ratio)

            # 图片位置 (居中偏右)
            left = (slide.slide_width - display_width) // 2 + Cm(2)
            top = (slide.slide_height - display_height) // 2

            # 重新读取图片数据
            image_data.seek(0)
            slide.shapes.add_picture(image_data, left, top, display_width, display_height)
            logger.info(f"Image inserted as shape at position ({left}, {top}), size: {display_width}x{display_height}")

        except Exception as e:
            logger.error(f"Failed to add image as shape: {str(e)}")

def ppt_generator_node(state: PPTState) -> PPTState:
    """
    PPT 渲染节点：最终生成 .pptx 文件并记录路径
    """
    logger.info("PPT Generator Node: Starting rendering...")
    logger.info(f"PPT Generator Node: Input state outline: {state.get('outline')}")
    logger.info(f"PPT Generator Node: Input state slides count: {len(state.get('slides', []))}")

    generator = PPTGenerator()

    try:
        file_path = generator.generate(state)
        logger.info(f"PPT Generator Node: Generated file: {file_path}")

        # 将生成的文件路径存入 state (需要更新状态定义以支持 file_path)
        # 或者我们暂时存入 error 作为传递手段，或者直接返回
        return {
            **state,
            "generated_file": file_path,
            "current_step": "completed"
        }
    except Exception as e:
        logger.exception("PPT Generator Node: Failed to render PPT.")
        return {
            **state,
            "error": f"Error in PPT generator: {str(e)}"
        }
